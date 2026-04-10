"""
Document processor: extracts text from PDF, DOCX, XLSX, PPTX, TXT, CSV files.
"""
from pathlib import Path


class DocumentProcessor:
    def extract_text(self, file_path: str) -> str:
        path = Path(file_path)
        ext = path.suffix.lower()

        extractors = {
            '.pdf': self._extract_pdf,
            '.doc': self._extract_docx,
            '.docx': self._extract_docx,
            '.xls': self._extract_xlsx,
            '.xlsx': self._extract_xlsx,
            '.ppt': self._extract_pptx,
            '.pptx': self._extract_pptx,
            '.txt': self._extract_txt,
            '.csv': self._extract_txt,
            '.md': self._extract_txt,
            '.png': self._extract_image,
            '.jpg': self._extract_image,
            '.jpeg': self._extract_image,
            '.gif': self._extract_image,
            '.webp': self._extract_image,
            '.bmp': self._extract_image,
            '.tiff': self._extract_image,
        }

        extractor = extractors.get(ext, self._extract_txt)
        try:
            return extractor(file_path)
        except Exception as e:
            return f"[Error extracting text from {path.name}: {e}]"

    def _extract_pdf(self, path: str) -> str:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)

                # Also extract tables
                tables = page.extract_tables()
                for table in tables:
                    for row in table:
                        if row:
                            text_parts.append(" | ".join(str(cell or "") for cell in row))
        return "\n\n".join(text_parts)

    def _extract_docx(self, path: str) -> str:
        from docx import Document
        doc = Document(path)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)

        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))
        return "\n\n".join(parts)

    def _extract_xlsx(self, path: str) -> str:
        from openpyxl import load_workbook
        wb = load_workbook(path, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"=== Hoja: {sheet_name} ===")
            for row in ws.iter_rows(values_only=True):
                cells = [str(cell) if cell is not None else "" for cell in row]
                if any(cells):
                    parts.append(" | ".join(cells))
        return "\n".join(parts)

    def _extract_pptx(self, path: str) -> str:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"--- Diapositiva {i} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            parts.append(para.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        parts.append(" | ".join(cells))
        return "\n\n".join(parts)

    def _extract_txt(self, path: str) -> str:
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()

    def _extract_image(self, path: str) -> str:
        """Extract text from image via OpenAI Vision API, fallback to placeholder."""
        import os, base64
        from pathlib import Path as _Path
        openai_key = os.environ.get("OPENAI_API_KEY", "")
        if not openai_key:
            return f"[Imagen adjunta: {_Path(path).name}]"
        try:
            from openai import OpenAI
            client = OpenAI(api_key=openai_key)
            with open(path, "rb") as f:
                b64 = base64.b64encode(f.read()).decode()
            ext = _Path(path).suffix.lower().lstrip('.')
            mime = {"jpg": "image/jpeg", "jpeg": "image/jpeg", "png": "image/png",
                    "gif": "image/gif", "webp": "image/webp", "bmp": "image/bmp"}.get(ext, "image/jpeg")
            resp = client.chat.completions.create(
                model="gpt-4o-mini", max_tokens=1024,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
                        {"type": "text", "text": (
                            "Extrae y transcribe todo el texto visible en esta imagen. "
                            "Si es una foto de apuntes o diapositivas, copia el contenido textual completo. "
                            "Si es un diagrama, describe su contenido. "
                            "Si no hay texto, describe brevemente la imagen."
                        )}
                    ]
                }]
            )
            return resp.choices[0].message.content or f"[Imagen: {_Path(path).name}]"
        except Exception:
            return f"[Imagen adjunta: {_Path(path).name}]"
